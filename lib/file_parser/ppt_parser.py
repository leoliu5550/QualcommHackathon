from pathlib import Path
from pptx import Presentation
from lib.file_parser.base_parser import BaseParser, ParseResult

class PptxParser(BaseParser):
    """PPTX 檔案解析器"""

    def parse(self, file_path: Path) -> ParseResult:
        try:
            presentation = Presentation(file_path)
            text_content = []
            char_count = 0

            for i, slide in enumerate(presentation.slides):
                slide_text = f"投影片 {i + 1}:\n"
                text_content.append(slide_text)
                char_count += len(slide_text)

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        shape_text = shape.text + '\n'
                        if char_count + len(shape_text) > self.char_limit:
                            remaining = self.char_limit - char_count
                            text_content.append(shape_text[:remaining])
                            char_count += remaining
                            break
                        text_content.append(shape_text)
                        char_count += len(shape_text)

                if char_count >= self.char_limit:
                    break

            content = ''.join(text_content)
            truncated_content, is_truncated = self._truncate_content(content)

            return ParseResult(
                success=True,
                content=truncated_content,
                file_type='pptx',
                original_length=len(content),
                truncated=is_truncated,
                file_path=str(file_path)
            )

        except Exception as e:
            return ParseResult(success=False, error=f"無法解析PPTX檔案: {str(e)}", file_path=str(file_path))
