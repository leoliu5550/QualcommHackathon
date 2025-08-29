"""PowerPoint presentation parser.

Extracts text from slides for content analysis.
"""

from pathlib import Path
from pptx import Presentation
from fileorg.parsers.base import BaseParser, ParseResult


class PptxParser(BaseParser):
    """Parser for PowerPoint presentations.
    
    Extracts text from all slides and shapes.
    Requires python-pptx library.
    """

    def parse(self, file_path: Path) -> ParseResult:
        """Extract text from PowerPoint presentation.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            ParseResult with slide text content
        """
        try:
            presentation = Presentation(file_path)
            text_content = []
            char_count = 0

            for i, slide in enumerate(presentation.slides):
                slide_text = f"Slide {i + 1}:\n"
                text_content.append(slide_text)
                char_count += len(slide_text)

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        shape_text = shape.text + "\n"
                        if char_count + len(shape_text) > self.char_limit:
                            remaining = self.char_limit - char_count
                            text_content.append(shape_text[:remaining])
                            char_count += remaining
                            break
                        text_content.append(shape_text)
                        char_count += len(shape_text)

                if char_count >= self.char_limit:
                    break

            content = "".join(text_content)
            truncated_content, is_truncated = self._truncate_content(content)

            return ParseResult(
                success=True,
                content=truncated_content,
                file_type="pptx",
                original_length=len(content),
                truncated=is_truncated,
                file_path=str(file_path),
            )

        except Exception as e:
            return ParseResult(
                success=False, error=f"Failed to parse PPTX: {str(e)}", file_path=str(file_path)
            )
