# ppt_parser.py
from pathlib import Path

from fileorg.file_ops.ports import IParser, ParserOutput


class PptParser(IParser):
    """Safe parser for PowerPoint presentations (.pptx only)."""

    def _truncate_content(self, content: str, char_limit: int) -> str:
        """Truncate content if it exceeds character limit."""
        if len(content) > char_limit:
            return content[:char_limit]
        return content

    def parse(self, file_path, char_limit: int) -> ParserOutput:
        """Parse a .pptx file and extract all text content."""
        try:
            from pptx import Presentation

            # 支援 str 或 Path
            file_path = Path(file_path)

            if file_path.suffix.lower() == ".ppt":
                return ParserOutput(success=False, content="", error="Unsupported file type: .ppt detected. Please convert to .pptx first.")

            if file_path.suffix.lower() != ".pptx":
                return ParserOutput(success=False, content="", error=f"Unsupported file type: {file_path.suffix}. Only .pptx is supported.")

            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slides_text.append(shape.text)

            full_text = "\n".join(slides_text)
            truncated_text = self._truncate_content(full_text, char_limit)

            return ParserOutput(success=True, content=truncated_text, error="")

        except Exception as e:
            return ParserOutput(success=False, content="", error=str(e))
